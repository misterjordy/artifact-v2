"""PPTX text extraction using python-pptx."""

import io

from pptx import Presentation

from artiFACT.modules.import_pipeline.extractors.base import ExtractorBase


class PptxExtractor(ExtractorBase):
    """Extract text from PPTX files."""

    def extract(self, content: bytes) -> str:
        prs = Presentation(io.BytesIO(content))
        parts: list[str] = []
        for slide in prs.slides:
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                parts.append("\n".join(slide_texts))
        return "\n\n".join(parts)
